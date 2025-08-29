#!/usr/bin/env python3
"""
PPTX Verification Script
Verifies that the generated PPTX contains the expected content from poster.html
"""

from pptx import Presentation
import os

def verify_pptx_content():
    """Verify the content of the generated PPTX file"""
    
    pptx_file = 'poster.pptx'
    
    if not os.path.exists(pptx_file):
        print(f"❌ Error: {pptx_file} not found!")
        return False
    
    print(f"✅ Found PPTX file: {pptx_file}")
    
    try:
        # Load the presentation
        prs = Presentation(pptx_file)
        
        print(f"📊 Presentation info:")
        print(f"   - Number of slides: {len(prs.slides)}")
        print(f"   - Slide width: {prs.slide_width}")
        print(f"   - Slide height: {prs.slide_height}")
        
        if len(prs.slides) == 0:
            print("❌ No slides found in presentation!")
            return False
        
        # Check the first (and only) slide
        slide = prs.slides[0]
        print(f"\n📄 Slide 1 content:")
        print(f"   - Number of shapes: {len(slide.shapes)}")
        
        # Extract all text content
        all_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        all_text.append(paragraph.text.strip())
        
        print(f"   - Text elements found: {len(all_text)}")
        
        # Check for key content elements
        key_elements = [
            "LLM4Hardware",
            "Generative AI for Chip Design",
            "github.com/FCHXWH823/LLM4Hardware",
            "Featured Submodules",
            "AutoChip",
            "Security Assertions", 
            "C2HLSC",
            "LLM4Verilog Generation",
            "LLM4Security",
            "LLM4C2HLS",
            "Research Projects"
        ]
        
        found_elements = []
        missing_elements = []
        
        for element in key_elements:
            found = False
            for text in all_text:
                if element in text:
                    found = True
                    found_elements.append(element)
                    break
            if not found:
                missing_elements.append(element)
        
        print(f"\n✅ Content verification:")
        print(f"   - Found elements: {len(found_elements)}/{len(key_elements)}")
        
        for element in found_elements:
            print(f"   ✓ {element}")
        
        if missing_elements:
            print(f"\n❌ Missing elements:")
            for element in missing_elements:
                print(f"   ✗ {element}")
        
        print(f"\n📝 All text content:")
        for i, text in enumerate(all_text, 1):
            print(f"   {i}. {text}")
        
        success = len(missing_elements) == 0
        if success:
            print(f"\n✅ Verification successful! PPTX contains all expected content.")
        else:
            print(f"\n⚠️  Verification partially successful. Some content may be missing.")
        
        return success
        
    except Exception as e:
        print(f"❌ Error reading PPTX file: {e}")
        return False

if __name__ == "__main__":
    verify_pptx_content()