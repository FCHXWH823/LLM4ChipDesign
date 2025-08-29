#!/usr/bin/env python3
"""
HTML to PPTX Converter for LLM4Hardware Poster
Converts poster.html to an equivalent PPTX presentation with one slide
"""

import os
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def create_poster_pptx():
    """Create a PPTX file from the poster.html content"""
    
    # Read the HTML file
    with open('poster.html', 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    # Parse HTML with BeautifulSoup
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Create new presentation
    prs = Presentation()
    
    # Use blank slide layout
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Add background gradient (approximating the CSS gradient)
    # Create a rectangle covering the entire slide
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(102, 126, 234)  # #667eea
    background.line.fill.background()
    
    # Add a white content area (simulating the poster white background)
    content_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.3),
        Inches(12.33), Inches(6.9)
    )
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_bg.line.color.rgb = RGBColor(200, 200, 200)
    content_bg.shadow.inherit = False
    
    # Header section with gradient background
    header_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.3),
        Inches(12.33), Inches(1.5)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(102, 126, 234)  # #667eea
    header_bg.line.fill.background()
    
    header_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.4), Inches(11.33), Inches(1.3)
    )
    header_frame = header_box.text_frame
    header_frame.clear()
    
    # Title
    title_p = header_frame.paragraphs[0]
    title_p.text = "🚀 LLM4Hardware"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_p = header_frame.add_paragraph()
    subtitle_p.text = "Generative AI for Chip Design"
    subtitle_p.font.size = Pt(20)
    subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # GitHub info section
    github_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.2), Inches(11.33), Inches(0.8)
    )
    github_frame = github_box.text_frame
    github_frame.clear()
    
    github_p = github_frame.paragraphs[0]
    github_p.text = "📦 Repository: github.com/FCHXWH823/LLM4Hardware"
    github_p.font.size = Pt(16)
    github_p.font.bold = True
    github_p.font.color.rgb = RGBColor(51, 51, 51)
    github_p.alignment = PP_ALIGN.CENTER
    
    # Featured projects section
    featured_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.2), Inches(12.33), Inches(1.8)
    )
    featured_frame = featured_box.text_frame
    featured_frame.clear()
    
    # Featured title
    featured_title = featured_frame.paragraphs[0]
    featured_title.text = "🌟 Featured Submodules"
    featured_title.font.size = Pt(20)
    featured_title.font.bold = True
    featured_title.font.color.rgb = RGBColor(229, 81, 0)  # #e65100
    featured_title.alignment = PP_ALIGN.CENTER
    
    # AutoChip project
    autochip_p = featured_frame.add_paragraph()
    autochip_p.text = "🔧 AutoChip: Generate functional Verilog modules from design prompts using LLMs with iterative error feedback"
    autochip_p.font.size = Pt(12)
    autochip_p.font.color.rgb = RGBColor(51, 51, 51)
    
    # Security project
    security_p = featured_frame.add_paragraph()
    security_p.text = "🛡️ Security Assertions: LLM-generated SystemVerilog assertions for hardware security verification"
    security_p.font.size = Pt(12)
    security_p.font.color.rgb = RGBColor(51, 51, 51)
    
    # C2HLSC project
    c2hlsc_p = featured_frame.add_paragraph()
    c2hlsc_p.text = "⚡ C2HLSC: Bridge software-to-hardware design gap using LLMs to refactor C code for HLS"
    c2hlsc_p.font.size = Pt(12)
    c2hlsc_p.font.color.rgb = RGBColor(51, 51, 51)
    
    # Main topics section (left column - Verilog)
    verilog_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.2), Inches(3.8), Inches(1.8)
    )
    verilog_frame = verilog_box.text_frame
    verilog_frame.clear()
    
    verilog_title = verilog_frame.paragraphs[0]
    verilog_title.text = "🔧 LLM4Verilog Generation"
    verilog_title.font.size = Pt(14)
    verilog_title.font.bold = True
    verilog_title.font.color.rgb = RGBColor(76, 175, 80)  # #4CAF50
    
    verilog_projects = [
        "• AutoChip: Functional Verilog generation with LLM feedback",
        "• VeriThoughts: Reasoning-based Verilog code generation", 
        "• ROME: Hierarchical prompting for complex chip design",
        "• PrefixLLM: LLM-aided prefix circuit design",
        "• Veritas: Deterministic Verilog synthesis from CNF"
    ]
    
    for project in verilog_projects:
        p = verilog_frame.add_paragraph()
        p.text = project
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(51, 51, 51)
    
    # Security section (middle column)
    security_box = slide.shapes.add_textbox(
        Inches(4.7), Inches(5.2), Inches(3.8), Inches(1.8)
    )
    security_frame = security_box.text_frame
    security_frame.clear()
    
    security_title = security_frame.paragraphs[0]
    security_title.text = "🛡️ LLM4Security"
    security_title.font.size = Pt(14)
    security_title.font.bold = True
    security_title.font.color.rgb = RGBColor(255, 87, 34)  # #FF5722
    
    security_projects = [
        "• Security Assertions: Hardware assertion generation for security",
        "• Hybrid-NL2SVA: Natural Language to SystemVerilog Assertion",
        "• OpenTitan RAG SVA: RAG system for OpenTitan SVA generation",
        "• LLMPirate: LLMs for black-box hardware IP piracy",
        "• FSM Testbench Gen: Testbench generation and bug detection"
    ]
    
    for project in security_projects:
        p = security_frame.add_paragraph()
        p.text = project
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(51, 51, 51)
    
    # C2HLS section (right column)
    c2hls_box = slide.shapes.add_textbox(
        Inches(8.9), Inches(5.2), Inches(3.8), Inches(1.8)
    )
    c2hls_frame = c2hls_box.text_frame
    c2hls_frame.clear()
    
    c2hls_title = c2hls_frame.paragraphs[0]
    c2hls_title.text = "⚡ LLM4C2HLS"
    c2hls_title.font.size = Pt(14)
    c2hls_title.font.bold = True
    c2hls_title.font.color.rgb = RGBColor(33, 150, 243)  # #2196F3
    
    c2hls_projects = [
        "• C2HLSC: Bridge software-to-hardware design gap",
        "• HLS Optimization: Automated C-to-HLS code transformation",
        "• Masala-CHAI: Large-scale SPICE netlist dataset"
    ]
    
    for project in c2hls_projects:
        p = c2hls_frame.add_paragraph()
        p.text = project
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(51, 51, 51)
    
    # Statistics section
    stats_box = slide.shapes.add_textbox(
        Inches(1), Inches(7.1), Inches(11.33), Inches(0.3)
    )
    stats_frame = stats_box.text_frame
    stats_frame.clear()
    
    stats_p = stats_frame.paragraphs[0]
    stats_p.text = "📊 12+ Research Projects • 7 Git Submodules • 10+ Published Papers • 3 Main Focus Areas"
    stats_p.font.size = Pt(12)
    stats_p.font.bold = True
    stats_p.font.color.rgb = RGBColor(255, 255, 255)
    stats_p.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    output_filename = 'poster.pptx'
    prs.save(output_filename)
    print(f"PPTX file saved as: {output_filename}")
    return output_filename

if __name__ == "__main__":
    create_poster_pptx()